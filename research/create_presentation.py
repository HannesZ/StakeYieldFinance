#!/usr/bin/env python3
"""Generate an introductory PowerPoint for StakeYieldFinance / StableYield Protocol."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGURES = os.path.join(os.path.dirname(__file__), "figures")
OUT = os.path.join(os.path.dirname(__file__), "StakeYieldFinance_Intro.pptx")

# Brand colours
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)      # deep navy
ACCENT  = RGBColor(0x63, 0x7D, 0xEA)       # muted indigo-blue
ACCENT2 = RGBColor(0x4E, 0xC9, 0xB0)       # teal/mint
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xC0, 0xC8, 0xE0)
DIMMED  = RGBColor(0x80, 0x8A, 0xA0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── helpers ──────────────────────────────────────────────────────────────────

def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
        pf = p._pPr
        from pptx.oxml.ns import qn
        buChar = pf.makeelement(qn('a:buChar'), {'char': '▸'})
        pf.append(buChar)
    return txBox

def add_image(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, width, height)
    return None

def accent_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def footer(slide, text="StakeYieldFinance — Confidential"):
    add_text(slide, Inches(0.5), H - Inches(0.5), Inches(6), Inches(0.4),
             text, size=9, color=DIMMED)

# ── Slide 1 — Title ─────────────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_bg(sl)
accent_bar(sl, Inches(0.8), Inches(2.6), Inches(4), Pt(4))
add_text(sl, Inches(0.8), Inches(2.8), Inches(11), Inches(1.2),
         "StakeYieldFinance", size=48, bold=True, color=WHITE)
add_text(sl, Inches(0.8), Inches(4.0), Inches(11), Inches(0.8),
         "Fixed-Rate Yield on Ethereum Staking — Powered by Actuarial Science",
         size=22, color=ACCENT2)
add_text(sl, Inches(0.8), Inches(5.2), Inches(6), Inches(0.5),
         "Hannes Zühlke  ·  May 2026", size=14, color=DIMMED)
footer(sl)

# ── Slide 2 — The Problem ───────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "The Problem", size=36, bold=True, color=ACCENT)
add_bullets(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5), [
    "Ethereum staking APR floats between ~2.5–4.5% — unpredictable",
    "Validators & LST holders bear full yield variance",
    "No native fixed-income instrument for staked ETH",
    "Existing protocols (Pendle, Notional) rely on AMM-implied rates — opaque pricing",
    "Institutions need predictable cash flows, not floating rewards",
], size=18)
img = add_image(sl, os.path.join(FIGURES, "08_implied_yield.png"),
                Inches(7.0), Inches(1.3), width=Inches(5.8))
add_text(sl, Inches(7.0), Inches(5.8), Inches(5.8), Inches(0.4),
         "Historical implied staking yield — volatile and hard to plan around",
         size=10, color=DIMMED, align=PP_ALIGN.CENTER)
footer(sl)

# ── Slide 3 — The Solution ──────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "The Solution: StableYield Protocol", size=36, bold=True, color=ACCENT)
add_bullets(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0), [
    "Deposit wstETH → receive a fixed-rate yield token (syLST)",
    "Locked rate at deposit time, redeemable at maturity",
    "Protocol reserve absorbs variance between floating & fixed yield",
    "Dynamic spread adjusts automatically based on reserve health",
    "Actuarial-grade risk pricing — transparent, on-chain solvency ratio",
    "No counterparty — the reserve IS the counterparty (insurance model)",
], size=18)
# simple diagram
add_text(sl, Inches(7.2), Inches(1.5), Inches(5.0), Inches(0.5),
         "How It Works", size=22, bold=True, color=ACCENT2)
steps = [
    ("①", "User deposits wstETH, selects maturity"),
    ("②", "Vault locks deposit, mints syLST at fixed rate"),
    ("③", "Floating yield accrues; surplus → reserve, deficit ← reserve"),
    ("④", "At maturity: redeem syLST for principal + fixed yield"),
]
y = Inches(2.2)
for num, desc in steps:
    add_text(sl, Inches(7.2), y, Inches(0.5), Inches(0.45), num, size=20, bold=True, color=ACCENT2)
    add_text(sl, Inches(7.8), y, Inches(4.8), Inches(0.45), desc, size=15, color=LIGHT)
    y += Inches(0.65)
footer(sl)

# ── Slide 4 — Yield Model Overview ──────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Yield Model — From Beacon Chain to Fixed Rates", size=36, bold=True, color=ACCENT)
add_bullets(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0), [
    "APR = f(total staked ETH) — deterministic issuance formula",
    "r = p × 5,256,000 / √B_total  (p ≈ 0.99 participation)",
    "Future APR depends on validator count trajectory",
    "Phase 1: Deterministic queue unwind (entry/exit queues → N_active)",
    "Phase 2: Merton jump-diffusion for stochastic arrivals/exits",
    "Monte Carlo → distribution of future APR paths",
    "Calibrated from 326k+ rows of beacon-chain history",
], size=17)
img = add_image(sl, os.path.join(FIGURES, "01_active_validators.png"),
                Inches(7.0), Inches(1.3), width=Inches(5.8))
add_text(sl, Inches(7.0), Inches(5.8), Inches(5.8), Inches(0.4),
         "Active validator count over time — the key state variable",
         size=10, color=DIMMED, align=PP_ALIGN.CENTER)
footer(sl)

# ── Slide 5 — Queue Dynamics ────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Queue Dynamics & Churn Limit", size=36, bold=True, color=ACCENT)
add_bullets(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.0), [
    "Entry/exit queues govern validator activation & withdrawal",
    "Churn limit: max(4, ⌊N/65536⌋) validators per epoch (capped at 8 post-Deneb)",
    "Queues create near-term deterministic APR trajectory",
    "Once queues clear, stochastic dynamics dominate",
    "Net flow currently exit-heavy → upward-sloping yield curve",
], size=17)
img = add_image(sl, os.path.join(FIGURES, "02_entry_exit_queues.png"),
                Inches(7.0), Inches(0.5), width=Inches(5.8))
img2 = add_image(sl, os.path.join(FIGURES, "05_net_flow.png"),
                 Inches(0.8), Inches(4.5), width=Inches(5.5))
footer(sl)

# ── Slide 6 — Yield Curve ───────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Constructing the Ethereum Staking Yield Curve", size=36, bold=True, color=ACCENT)
add_bullets(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), [
    "Monte Carlo simulation → 10,000 APR paths",
    "Discount factors: D(T) = E[exp(-∫₀ᵀ r(t) dt)]",
    "Zero-coupon rates extracted from discount factors",
    "Fan chart shows uncertainty bands (p5–p95)",
    "Foundation for pricing fixed-rate instruments",
], size=17)
img = add_image(sl, os.path.join(FIGURES, "09_yield_curve_recalibrated.png"),
                Inches(7.0), Inches(0.3), width=Inches(5.8))
img2 = add_image(sl, os.path.join(FIGURES, "10_apr_fan_chart.png"),
                 Inches(0.8), Inches(4.2), width=Inches(5.5))
footer(sl)

# ── Slide 7 — Protocol Architecture ─────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Protocol Architecture", size=36, bold=True, color=ACCENT)

# Architecture boxes
boxes = [
    ("StableYield Vault", "Main entry: deposit, redeem, harvest\nHolds all escrowed wstETH", Inches(4.5), Inches(1.6)),
    ("syLST (ERC-1155)", "Fixed-yield tokens\nOne ID per maturity series", Inches(1.0), Inches(3.6)),
    ("Reserve Manager", "Solvency tracking\nSurplus/deficit management", Inches(4.5), Inches(3.6)),
    ("Spread Calculator", "Dynamic spread from κ\nPure math, no state", Inches(8.0), Inches(3.6)),
    ("SYLD Token (ERC-20)", "Governance + backstop\nResidual claim on reserves", Inches(4.5), Inches(5.4)),
]
for title, desc, x, y in boxes:
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.2), Inches(1.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x24, 0x3B)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT2
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

# Frontend box
shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.6), Inches(3.2), Inches(1.4))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x24, 0x3B)
shape.line.color.rgb = ACCENT2
shape.line.width = Pt(1.5)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Frontend (Swap GUI)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT2
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Rate display, deposit flow\nHTML/JS interface"
p2.font.size = Pt(11)
p2.font.color.rgb = LIGHT
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER

footer(sl)

# ── Slide 8 — Dynamic Spread ────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Dynamic Spread — Self-Stabilising Risk Pricing", size=36, bold=True, color=ACCENT)

add_text(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "s(κ) = s_base × (1 + α × max(0, κ_target/κ − 1)^β)",
         size=16, bold=True, color=ACCENT2, font_name="Consolas")

add_bullets(sl, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), [
    "κ = Reserve / Liabilities (solvency ratio)",
    "κ ≥ 1.5 (target): spread = 50bp — attractive rates",
    "κ = 0.7: spread ≈ 115bp — gentle deterrent",
    "κ = 0.5: spread ≈ 470bp — strong deterrent",
    "κ = 0.3 (critical): spread ≈ 5,050bp — effectively closed",
    "No governance intervention needed — pure math",
    "Smooth transition, no cliff edges",
], size=16)

# Table of parameters
params = [
    ("s_base", "50 bp", "Base spread"),
    ("α", "10", "Stress multiplier"),
    ("β", "3", "Convexity (cubic)"),
    ("κ_target", "1.5", "Target solvency"),
    ("κ_critical", "0.3", "Emergency threshold"),
]
ty = Inches(1.5)
add_text(sl, Inches(7.5), ty, Inches(5.0), Inches(0.4),
         "Key Parameters", size=18, bold=True, color=ACCENT2)
ty += Inches(0.5)
for param, val, desc in params:
    add_text(sl, Inches(7.5), ty, Inches(1.5), Inches(0.35), param, size=14, bold=True, color=ACCENT2, font_name="Consolas")
    add_text(sl, Inches(9.0), ty, Inches(1.0), Inches(0.35), val, size=14, color=WHITE)
    add_text(sl, Inches(10.2), ty, Inches(2.5), Inches(0.35), desc, size=13, color=DIMMED)
    ty += Inches(0.38)

footer(sl)

# ── Slide 9 — Reserve & Solvency ────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Reserve Management — Actuarial Framework", size=36, bold=True, color=ACCENT)

add_bullets(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), [
    "Reserve absorbs variance: surplus → reserve, deficit ← reserve",
    "Inspired by Swiss Solvency Test (SST) / Solvency II",
    "Explicit SCR (Solvency Capital Requirement) from Monte Carlo",
    "SYLD token = subordinated backstop (MakerDAO MKR model)",
    "syLST holders are senior — claims always honoured first",
], size=17)

add_text(sl, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.5),
         "Credit Hierarchy", size=20, bold=True, color=ACCENT2)
hierarchy = [
    ("SENIOR", "syLST Depositors — fixed yield guaranteed by reserve", ACCENT2),
    ("JUNIOR", "SYLD Holders — residual claim, bear tail risk via dilution", ACCENT),
]
hy = Inches(4.6)
for label, desc, col in hierarchy:
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), hy, Inches(5.5), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x24, 0x3B)
    shape.line.color.rgb = col
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"  {label}:  {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT
    p.font.name = "Calibri"
    hy += Inches(0.65)

# Right side — reserve sizing
add_text(sl, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
         "Initial Reserve Sizing (10k wstETH TVL)", size=18, bold=True, color=ACCENT2)
reserve_items = [
    "Annual yield variance: ~7bp (from model)",
    "Tail risk (p1 scenario): ~15bp below expected",
    "1-year SCR: ~15 wstETH (0.15% of TVL)",
    "Target reserve (150%): ~22.5 wstETH",
    "Conservative with buffers: 3–5% of TVL",
    "Remarkably capital-efficient due to low model uncertainty",
]
add_bullets(sl, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.0), reserve_items, size=15)
footer(sl)

# ── Slide 10 — Token Design ─────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Token Design", size=36, bold=True, color=ACCENT)

# syLST
add_text(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "syLST — Stable Yield LST", size=24, bold=True, color=ACCENT2)
add_bullets(sl, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.0), [
    "ERC-1155: one token ID per maturity series (Q3-26, Q4-26, …)",
    "Claim to (1 + r_fixed × T/365) wstETH at maturity",
    "Fully transferable — tradeable on DEXes",
    "Composable: usable as collateral in lending protocols",
    "No early exit from vault — sell on secondary market",
], size=16)

# SYLD
add_text(sl, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
         "SYLD — Governance Token", size=24, bold=True, color=ACCENT2)
add_bullets(sl, Inches(7.0), Inches(2.2), Inches(5.5), Inches(3.0), [
    "ERC-20 with fixed supply + emergency mint",
    "Governance votes on protocol parameters",
    "Stake to receive excess reserve distributions",
    "Emergency backstop: mint & auction to recapitalize",
    "Subordinated to syLST (bears tail risk)",
], size=16)

# Comparison table
add_text(sl, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
         "vs. Existing Protocols", size=20, bold=True, color=ACCENT2)
add_text(sl, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5),
         "StableYield: model + actuarial pricing  ·  Pendle: AMM-implied  ·  Notional: pool-based\n"
         "Only protocol with explicit actuarial reserve & transparent on-chain solvency ratio",
         size=14, color=LIGHT)
footer(sl)

# ── Slide 11 — Discount Factors & Swap Pricer ───────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "From Yield Curves to Products", size=36, bold=True, color=ACCENT)
add_bullets(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), [
    "Discount factors from MC simulation → zero-coupon rates",
    "Fixed-rate quotes derived directly from the term structure",
    "Interest rate swap pricing: pay fixed, receive floating staking yield",
    "Interactive swap pricer GUI built for exploration",
], size=17)
img = add_image(sl, os.path.join(FIGURES, "11_discount_factors.png"),
                Inches(7.0), Inches(0.3), width=Inches(5.8))
img2 = add_image(sl, os.path.join(FIGURES, "swap_pricer_screenshot.png"),
                 Inches(0.8), Inches(4.0), width=Inches(5.5))
footer(sl)

# ── Slide 12 — Roadmap ──────────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Roadmap", size=36, bold=True, color=ACCENT)

phases = [
    ("Phase 1 — Testnet", "Core contracts, single maturity, Hardhat tests, GUI integration", ACCENT2),
    ("Phase 2 — Audit & Beta", "Professional audit, Holesky deployment, bug bounty, multi-series", ACCENT),
    ("Phase 3 — Mainnet", "Conservative launch, initial reserve, SYLD token, short tenors (≤6m)", RGBColor(0xE0, 0xA0, 0x50)),
    ("Phase 4 — DeFi Integration", "On-chain yield model, syLST as collateral, L2 deployment, multi-LST", RGBColor(0xC0, 0x60, 0xD0)),
]
y = Inches(1.6)
for title, desc, col in phases:
    # Phase box
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y, Inches(10.5), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x24, 0x3B)
    shape.line.color.rgb = col
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = col
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT
    p2.font.name = "Calibri"
    # dot
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.3), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    y += Inches(1.35)

footer(sl)

# ── Slide 13 — Key Differentiators ──────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         "Why StakeYieldFinance?", size=36, bold=True, color=ACCENT)

diffs = [
    ("🔬", "Model-Driven Pricing", "Actuarial yield model, not AMM guesswork — transparent, verifiable rates"),
    ("🛡️", "Explicit Risk Management", "On-chain solvency ratio, SCR, dynamic spread — SST/Solvency II inspired"),
    ("💰", "Capital Efficient", "Reserve requirement ~3–5% of TVL — low model uncertainty = small buffer"),
    ("🔗", "DeFi Composable", "syLST tradeable, usable as collateral — yield curve arbitrage opportunities"),
    ("⚖️", "Clear Credit Hierarchy", "syLST senior, SYLD junior — well-understood waterfall structure"),
    ("🤖", "Self-Stabilising", "Dynamic spread adjusts automatically — no governance fire drills"),
]
y = Inches(1.5)
for emoji, title, desc in diffs:
    add_text(sl, Inches(0.8), y, Inches(0.6), Inches(0.5), emoji, size=24)
    add_text(sl, Inches(1.5), y, Inches(3.5), Inches(0.5), title, size=18, bold=True, color=ACCENT2)
    add_text(sl, Inches(5.2), y, Inches(7.5), Inches(0.5), desc, size=15, color=LIGHT)
    y += Inches(0.75)
footer(sl)

# ── Slide 14 — Thank You / Contact ──────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
accent_bar(sl, Inches(5.0), Inches(2.8), Inches(3.5), Pt(4))
add_text(sl, Inches(0.8), Inches(2.9), Inches(11.5), Inches(1.0),
         "Thank You", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6),
         "Hannes Zühlke  ·  johanneszuehlke@fastmail.com  ·  github.com/HannesZ",
         size=16, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.6),
         "StakeYieldFinance — Fixed-Rate Yield Meets Actuarial Precision",
         size=18, color=ACCENT2, align=PP_ALIGN.CENTER)
footer(sl)

# ── Save ─────────────────────────────────────────────────────────────────────

prs.save(OUT)
print(f"✅ Saved → {OUT}")
